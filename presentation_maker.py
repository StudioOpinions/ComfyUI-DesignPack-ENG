import json
import os
import torch
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import io
import glob

class PresentationMaker:
    @classmethod
    def INPUT_TYPES(cls):
        # Retrieve font files
        font_folder = os.path.join(os.path.dirname(__file__), "fonts")
        font_files = []
        if os.path.exists(font_folder):
            font_files = [f for f in os.listdir(font_folder) if f.endswith('.ttf')]
        
        if not font_files:
            font_files = ["default"]
        
        return {
            "required": {
                "images": ("IMAGE",),
                "json_data": ("STRING",),  # ensure plain string input
                "template": (["Classic Minimal", "Modern Business", "Cozy Home", "Stylish Elegance", "Fresh Nature"], {"default": "Classic Minimal"}),
                "font": (font_files, {"default": font_files[0] if font_files else "default"}),
                "output_path": ("STRING", {"default": "./output/presentation.pptx"}),
            }
        }
    
    INPUT_IS_LIST = (True, False, False, False, False)  # only images is a list
    RETURN_TYPES = ("STRING",)
    RETURN_NAMES = ("result",)
    
    FUNCTION = "make_presentation"
    CATEGORY = "presentation"
    
    def __init__(self):
        # Template configurations
        self.templates = {
            "Classic Minimal": {
                "title_color": (44, 62, 80),
                "content_color": (52, 73, 94),
                "bg_color": (248, 249, 250),
                "accent_color": (52, 152, 219)
            },
            "Modern Business": {
                "title_color": (45, 52, 54),
                "content_color": (99, 110, 114),
                "bg_color": (255, 255, 255),
                "accent_color": (0, 123, 255)
            },
            "Cozy Home": {
                "title_color": (139, 69, 19),
                "content_color": (101, 67, 33),
                "bg_color": (255, 248, 220),
                "accent_color": (205, 133, 63)
            },
            "Stylish Elegance": {
                "title_color": (75, 0, 130),
                "content_color": (105, 105, 105),
                "bg_color": (245, 245, 245),
                "accent_color": (138, 43, 226)
            },
            "Fresh Nature": {
                "title_color": (34, 139, 34),
                "content_color": (85, 107, 47),
                "bg_color": (240, 255, 240),
                "accent_color": (60, 179, 113)
            }
        }
    
    def tensor_to_pil(self, tensor):
        """Convert tensor to a PIL image with memory optimization"""
        try:
            if isinstance(tensor, torch.Tensor):
                if tensor.dim() == 4:
                    tensor = tensor.squeeze(0)
                
                # Check image size and downscale if necessary
                height, width = tensor.shape[:2]
                max_dimension = 2048  # limit maximum dimension
                
                if height > max_dimension or width > max_dimension:
                    scale = max_dimension / max(height, width)
                    new_height = int(height * scale)
                    new_width = int(width * scale)
                    
                    # Use torch interpolation to resize, more memory efficient than PIL
                    tensor = torch.nn.functional.interpolate(
                        tensor.unsqueeze(0).permute(0, 3, 1, 2),  # [1, C, H, W]
                        size=(new_height, new_width),
                        mode='bilinear',
                        align_corners=False
                    ).squeeze(0).permute(1, 2, 0)  # [H, W, C]
                
                # Convert to numpy and clamp to valid range
                image_np = (tensor.cpu().numpy() * 255).astype(np.uint8)
                return Image.fromarray(image_np)
            return tensor
        except Exception as e:
            print(f"Tensor conversion error: {e}")
            # Create a default small image
            default_img = Image.new('RGB', (800, 600), color='white')
            return default_img
    
    def resize_image_to_fit(self, image, max_width, max_height, maintain_aspect=True):
        """Resize image to fit target dimensions with memory optimization"""
        try:
            # Skip processing if the image is already small
            if image.width <= max_width and image.height <= max_height and maintain_aspect:
                return image
            
            # Compute new size
            if maintain_aspect:
                # Maintain aspect ratio
                ratio = min(max_width / image.width, max_height / image.height)
                new_width = int(image.width * ratio)
                new_height = int(image.height * ratio)
                
                # Pre-scale if ratio is very small to avoid memory issues
                if ratio < 0.5:
                    # Step scaling: first resize to an intermediate size
                    intermediate_ratio = 0.7
                    intermediate_width = int(image.width * intermediate_ratio)
                    intermediate_height = int(image.height * intermediate_ratio)
                    image = image.resize((intermediate_width, intermediate_height), Image.Resampling.LANCZOS)
                    
                    # Recalculate final ratio
                    ratio = min(max_width / image.width, max_height / image.height)
                    new_width = int(image.width * ratio)
                    new_height = int(image.height * ratio)
                
            else:
                # 16:9 full-bleed crop
                target_ratio = 16/9
                img_ratio = image.width / image.height
                
                # Check if significant scaling is needed
                scale_factor = min(max_width / image.width, max_height / image.height)
                if scale_factor < 0.5:
                    # Pre-scale to near target size
                    temp_width = int(max_width * 1.2)
                    temp_height = int(max_height * 1.2)
                    image = image.resize((temp_width, temp_height), Image.Resampling.LANCZOS)
                
                if img_ratio > target_ratio:
                    # Image too wide, crop width
                    new_height = image.height
                    new_width = int(new_height * target_ratio)
                    x = (image.width - new_width) // 2
                    image = image.crop((x, 0, x + new_width, new_height))
                else:
                    # Image too tall, crop height
                    new_width = image.width
                    new_height = int(new_width / target_ratio)
                    y = (image.height - new_height) // 2
                    image = image.crop((0, y, new_width, y + new_height))
                
                new_width = max_width
                new_height = max_height
            
            # Final resize
            return image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
        except MemoryError:
            # Degraded processing when memory is low
            print(f"Low memory, attempting degraded processing for image {image.width}x{image.height}")

            # Use a more aggressive scaling strategy
            if maintain_aspect:
                ratio = min(max_width / image.width, max_height / image.height)
                # Further reduce ratio to save memory
                safe_ratio = min(ratio * 0.8, 0.5)
                new_width = int(image.width * safe_ratio)
                new_height = int(image.height * safe_ratio)
            else:
                # Directly set to a smaller target size
                new_width = int(max_width * 0.8)
                new_height = int(max_height * 0.8)

            # Use lower quality scaling to conserve memory
            return image.resize((new_width, new_height), Image.Resampling.NEAREST)
            
        except Exception as e:
            print(f"Image processing error: {e}")
            # Return a small version of the original image
            return image.resize((min(800, image.width), min(600, image.height)), Image.Resampling.NEAREST)
    
    def add_image_to_slide(self, slide, image, left, top, width, height):
        """Add an image to a slide"""
        img_stream = io.BytesIO()
        image.save(img_stream, format='PNG')
        img_stream.seek(0)
        
        pic = slide.shapes.add_picture(img_stream, left, top, width, height)
        return pic
    
    def generate_unique_filename(self, output_path):
        """Generate a unique filename to avoid overwriting"""
        if not os.path.exists(output_path):
            return output_path
        
        # Split file path, name, and extension
        dir_path = os.path.dirname(output_path)
        filename = os.path.basename(output_path)
        name, ext = os.path.splitext(filename)
        
        # Find an available filename
        counter = 1
        while True:
            new_filename = f"{name}{counter:02d}{ext}"
            new_path = os.path.join(dir_path, new_filename)
            if not os.path.exists(new_path):
                return new_path
            counter += 1
            
            # Prevent infinite loops; try at most 100 times
            if counter > 100:
                import time
                timestamp = int(time.time())
                new_filename = f"{name}_{timestamp}{ext}"
                return os.path.join(dir_path, new_filename)
    
    def make_presentation(self, images, json_data, template, font, output_path):
        try:
            # Per ComfyUI, all parameters are lists when INPUT_IS_LIST=True
            # Extract actual values from lists
            json_string = json_data[0] if isinstance(json_data, list) and json_data else "{}"
            template_name = template[0] if isinstance(template, list) and template else "Classic Minimal"
            font_name = font[0] if isinstance(font, list) and font else "default"
            output_file = output_path[0] if isinstance(output_path, list) and output_path else "./presentation.pptx"
            
            # Process image list - flatten batches
            all_images = []
            for img_batch in images:  # images is a list of image batches
                if isinstance(img_batch, torch.Tensor):
                    # If batch tensor, split each image
                    if img_batch.dim() == 4:  # [batch, height, width, channels]
                        for i in range(img_batch.shape[0]):
                            all_images.append(img_batch[i])
                    else:  # single image
                        all_images.append(img_batch)
                else:
                    all_images.append(img_batch)
            
            # Parse JSON data
            try:
                presentation_data = json.loads(json_string)
            except json.JSONDecodeError as e:
                return (f"JSON parse error: {str(e)}; received data: {str(json_string)[:200]}",)
            
            # Verify image count matches JSON data
            if len(all_images) != len(presentation_data):
                return (f"Image count ({len(all_images)}) does not match JSON entries ({len(presentation_data)})",)
            
            # Create new presentation
            prs = Presentation()
            
            # Set presentation size to 16:9
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Retrieve template configuration
            template_config = self.templates.get(template_name, self.templates["Classic Minimal"])
            
            # Process each page
            for i, (image, data_item) in enumerate(zip(all_images, presentation_data)):
                print(f"Processing page {i+1}...")

                # Extract data
                value = data_item.get("value", "")
                description = data_item.get("description", "")
                
                # Check if this is the second page (proposal style overview)
                is_second_page = (value == "Proposal Style Overview" and
                                "topic1" in data_item and
                                "summary1" in data_item)
                
                # Skip empty values
                if not value:
                    continue
                
                # Add new slide
                slide_layout = prs.slide_layouts[6]  # blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*template_config["bg_color"])
                
                # Convert image
                print("Converting image...")
                pil_image = self.tensor_to_pil(image)
                print(f"Image size: {pil_image.width}x{pil_image.height}")
                
                # First and last pages - full-bleed design
                if i == 0 or i == len(all_images) - 1:
                    print("Scaling to full bleed...")
                    # Full-bleed 16:9 - use smaller size to save memory
                    slide_width_px = 1920  # lower resolution
                    slide_height_px = 1080
                    
                    processed_image = self.resize_image_to_fit(
                        pil_image, 
                        slide_width_px, 
                        slide_height_px, 
                        maintain_aspect=False
                    )
                    print(f"Scaled size: {processed_image.width}x{processed_image.height}")

                    # Add background image
                    self.add_image_to_slide(
                        slide, processed_image, 
                        0, 0, 
                        prs.slide_width, prs.slide_height
                    )
                    
                    # Use layered approach to create semi-transparent background
                    # Position and size of text area
                    text_left = Inches(1)
                    text_top = Inches(2.5)
                    text_width = Inches(11.33)
                    text_height = Inches(2.5)
                    
                    # ======= Bottom layer: semi-transparent background box (simulate dark) =======
                    from pptx.enum.shapes import MSO_SHAPE
                    bg_box = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        text_left,
                        text_top,
                        text_width,
                        text_height
                    )
                    bg_box.fill.solid()
                    bg_box.fill.fore_color.rgb = RGBColor(40, 40, 40)  # simulate 70% black transparency
                    bg_box.line.fill.background()  # no border
                    
                    # ======= Top layer: text box with transparent background =======
                    text_box = slide.shapes.add_textbox(
                        text_left, text_top, text_width, text_height
                    )
                    text_frame = text_box.text_frame
                    text_frame.clear()
                    text_frame.margin_left = Inches(0.2)
                    text_frame.margin_right = Inches(0.2)
                    text_frame.margin_top = Inches(0.2)
                    text_frame.margin_bottom = Inches(0.2)
                    
                    # Title text
                    title_p = text_frame.add_paragraph()
                    title_p.text = str(value)
                    title_p.alignment = PP_ALIGN.CENTER
                    title_p.font.size = Pt(44)
                    title_p.font.color.rgb = RGBColor(255, 255, 255)
                    title_p.font.bold = True
                    
                    # Subtitle
                    if description:
                        subtitle_p = text_frame.add_paragraph()
                        subtitle_p.text = str(description)
                        subtitle_p.alignment = PP_ALIGN.CENTER
                        subtitle_p.font.size = Pt(24)
                        subtitle_p.font.color.rgb = RGBColor(240, 240, 240)
                    
                    # Make text box background transparent
                    text_box.fill.background()  # transparent background
                    text_box.line.fill.background()  # no border
                
                elif is_second_page:
                    # Second page - proposal style overview - special layout
                    print("Processing special layout for page 2...")

                    # Add full background image
                    slide_width_px = 1920
                    slide_height_px = 1080
                    
                    processed_image = self.resize_image_to_fit(
                        pil_image, 
                        slide_width_px, 
                        slide_height_px, 
                        maintain_aspect=False
                    )
                    
                    # Add background image
                    self.add_image_to_slide(
                        slide, processed_image, 
                        0, 0, 
                        prs.slide_width, prs.slide_height
                    )
                    
                    # Add title (top) - same style as page 3
                    title_box = slide.shapes.add_textbox(
                        Inches(1), Inches(0.5),
                        Inches(11.33), Inches(1)
                    )
                    title_frame = title_box.text_frame
                    title_frame.clear()
                    title_frame.margin_left = Inches(0.2)
                    title_frame.margin_right = Inches(0.2)
                    title_frame.margin_top = Inches(0.2)
                    title_frame.margin_bottom = Inches(0.2)
                    
                    title_p = title_frame.add_paragraph()
                    title_p.text = str(value)
                    title_p.alignment = PP_ALIGN.CENTER
                    title_p.font.size = Pt(36)
                    title_p.font.color.rgb = RGBColor(*template_config["title_color"])
                    title_p.font.bold = True
                    
                    # Make text box background transparent
                    title_box.fill.background()
                    title_box.line.fill.background()
                    
                    # Calculate positions for four text boxes (2x2 layout) for better centering
                    box_width = Inches(5.4)
                    box_height = Inches(2.2)
                    margin_x = Inches(0.3)  # reduce horizontal margin
                    margin_y = Inches(0.3)
                    start_y = Inches(2.2)
                    
                    # Compute total width for centering
                    total_width = 2 * box_width + margin_x  # width of two boxes plus gap
                    start_x = (prs.slide_width - total_width) / 2  # center start position
                    
                    topics = [
                        data_item.get("topic1", ""),
                        data_item.get("topic2", ""),
                        data_item.get("topic3", ""),
                        data_item.get("topic4", "")
                    ]
                    
                    summaries = [
                        data_item.get("summary1", ""),
                        data_item.get("summary2", ""),
                        data_item.get("summary3", ""),
                        data_item.get("summary4", "")
                    ]
                    
                    # Create four rounded text boxes
                    for idx, (topic, summary) in enumerate(zip(topics, summaries)):
                        row = idx // 2
                        col = idx % 2
                        
                        box_left = start_x + col * (box_width + margin_x)
                        box_top = start_y + row * (box_height + margin_y)
                        
                        # Create plain rectangular background box (no rounded corners)
                        bg_box = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            box_left, box_top,
                            box_width, box_height
                        )
                        bg_box.fill.solid()
                        bg_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # white background
                        bg_box.line.color.rgb = RGBColor(*template_config["accent_color"])
                        bg_box.line.width = Pt(2)
                        
                        # Add title text
                        title_text_box = slide.shapes.add_textbox(
                            box_left + Inches(0.2), box_top + Inches(0.2),
                            box_width - Inches(0.4), Inches(0.6)
                        )
                        title_text_frame = title_text_box.text_frame
                        title_text_frame.clear()
                        title_text_frame.margin_left = Inches(0.1)
                        title_text_frame.margin_right = Inches(0.1)
                        title_text_frame.margin_top = Inches(0.05)
                        title_text_frame.margin_bottom = Inches(0.05)
                        
                        title_text_p = title_text_frame.add_paragraph()
                        title_text_p.text = str(topic)
                        title_text_p.alignment = PP_ALIGN.CENTER
                        title_text_p.font.size = Pt(18)
                        title_text_p.font.color.rgb = RGBColor(*template_config["title_color"])
                        title_text_p.font.bold = True
                        
                        title_text_box.fill.background()
                        title_text_box.line.fill.background()
                        
                        # Remove connector lines and add content text directly (closer to title)
                        content_text_box = slide.shapes.add_textbox(
                            box_left + Inches(0.2), box_top + Inches(0.9),
                            box_width - Inches(0.4), box_height - Inches(1.1)
                        )
                        content_text_frame = content_text_box.text_frame
                        content_text_frame.clear()
                        content_text_frame.margin_left = Inches(0.1)
                        content_text_frame.margin_right = Inches(0.1)
                        content_text_frame.margin_top = Inches(0.05)
                        content_text_frame.margin_bottom = Inches(0.05)
                        content_text_frame.word_wrap = True
                        
                        content_text_p = content_text_frame.add_paragraph()
                        content_text_p.text = str(summary)
                        content_text_p.alignment = PP_ALIGN.LEFT
                        content_text_p.font.size = Pt(14)
                        content_text_p.font.color.rgb = RGBColor(*template_config["content_color"])
                        content_text_p.line_spacing = 1.2
                        
                        content_text_box.fill.background()
                        content_text_box.line.fill.background()
                
                
                else:
                    print(f"Processing content page {i+1}, keeping original image size...")
                    # Content page - side-by-side layout, keep original size
                    # No scaling; use original image directly
                    processed_image = pil_image

                    # Calculate image display size in presentation
                    display_width = Inches(5.5)
                    display_height = Inches(4.6)
                    
                    # If the original image is too large, slightly scale it to fit
                    if pil_image.width > 1200 or pil_image.height > 1000:
                        print(f"Image too large ({pil_image.width}x{pil_image.height}), applying moderate scaling...")
                        # Only scale when the image is obviously oversized
                        processed_image = self.resize_image_to_fit(
                            pil_image, 1200, 1000, maintain_aspect=True
                        )
                        print(f"Scaled size: {processed_image.width}x{processed_image.height}")
                    else:
                        print(f"Image size acceptable, keeping original: {pil_image.width}x{pil_image.height}")

                    # Add image (right side)
                    print("Adding image to presentation...")
                    self.add_image_to_slide(
                        slide, processed_image,
                        Inches(7.3), Inches(1.5),
                        display_width, display_height
                    )
                    
                    # Add title (top left)
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5),
                        Inches(6.5), Inches(1)
                    )
                    title_frame = title_box.text_frame
                    title_p = title_frame.add_paragraph()
                    title_p.text = str(value)
                    title_p.font.size = Pt(32)
                    title_p.font.color.rgb = RGBColor(*template_config["title_color"])
                    title_p.font.bold = True
                    
                    # Add content (bottom left)
                    content_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.6),
                        Inches(6.5), Inches(4.5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True
                    
                    # Increase font size to reduce whitespace and add paragraph breaks
                    if description:
                        # Split description into sentences and add paragraph breaks
                        description_text = str(description)
                        sentences = description_text.split('.')

                        for idx, sentence in enumerate(sentences):
                            if sentence.strip():  # Ignore empty strings
                                content_p = content_frame.add_paragraph()
                                content_p.text = sentence.strip() + ('.' if idx < len(sentences) - 1 and sentence.strip() else '')
                                content_p.font.size = Pt(17)  # increase from 16 to 17
                                content_p.font.color.rgb = RGBColor(*template_config["content_color"])
                                content_p.line_spacing = 1.3

                                # paragraph spacing
                                if idx < len(sentences) - 1:
                                    content_p.space_after = Pt(6)
                    else:
                        # Fallback empty content
                        content_p = content_frame.add_paragraph()
                        content_p.text = ""
                        content_p.font.size = Pt(17)
                        content_p.font.color.rgb = RGBColor(*template_config["content_color"])
                        content_p.line_spacing = 1.4
            
            # Ensure output directory exists
            output_dir = os.path.dirname(output_file)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # Generate unique filename to avoid overwrite
            unique_output_path = self.generate_unique_filename(output_file)
            
            # Save presentation
            prs.save(unique_output_path)
            
            return (f"Presentation generated successfully: {unique_output_path}",)
            
        except Exception as e:
            error_msg = f"Error generating presentation: {str(e)}"
            print(f"Detailed error: {e}")  # for debugging
            import traceback
            print(f"Traceback: {traceback.format_exc()}")
            return (error_msg,)

NODE_CLASS_MAPPINGS = {
    "PresentationMaker": PresentationMaker
}

NODE_DISPLAY_NAME_MAPPINGS = {
    "PresentationMaker": "Presentation Maker"
}